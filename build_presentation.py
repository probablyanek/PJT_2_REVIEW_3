import os
import subprocess
import shutil
import sys
from pptx import Presentation
from pptx.util import Inches

def run_command(command, description):
    """Utility to run shell commands and handle output/errors."""
    print(f"--- {description} ---")
    try:
        # Use shell=True for compatibility with Windows command strings
        result = subprocess.run(command, shell=True, capture_output=True, text=True)
        if result.returncode != 0:
            # pdflatex often returns 1 for minor Beamer/TikZ warnings; we only log it
            if "pdflatex" in command:
                print(f"Warning: pdflatex returned code {result.returncode} (common with Beamer). Continuing...")
            else:
                print(f"Error during {description}:")
                print(result.stderr)
                return False
        return True
    except Exception as e:
        print(f"Execution failed: {e}")
        return False

def build_presentation(tex_file="main.tex", output_name="main", dpi=400):
    """Main subroutine to go from .tex to .pdf to .pptx"""
    pdf_file = f"{output_name}.pdf"
    pptx_file = f"{output_name}.pptx"
    temp_dir = "temp_build_assets"

    # 1. Compile LaTeX (Double pass for references and TikZ)
    success = run_command(f"pdflatex -interaction=nonstopmode {tex_file}", "LaTeX Compilation (Pass 1)")
    if not success: return
    
    run_command(f"pdflatex -interaction=nonstopmode {tex_file}", "LaTeX Compilation (Pass 2)")

    if not os.path.exists(pdf_file):
        print(f"Error: {pdf_file} was not found. Compilation likely failed.")
        return

    # 2. Convert PDF to high-DPI images using pdftoppm
    if os.path.exists(temp_dir):
        shutil.rmtree(temp_dir)
    os.makedirs(temp_dir)

    print(f"Extracting images at {dpi} DPI...")
    success = run_command(f"pdftoppm -png -r {dpi} {pdf_file} {temp_dir}/slide", "PDF to Image Conversion")
    if not success: return

    # 3. Assemble the PowerPoint Presentation
    print(f"Assembling {pptx_file}...")
    prs = Presentation()
    
    # 16:9 Aspect Ratio (13.333 x 7.5 inches)
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    # Get images and sort them numerically (slide-1.png, slide-2.png, etc.)
    images = [f for f in os.listdir(temp_dir) if f.endswith('.png')]
    images.sort(key=lambda x: int(x.split('-')[-1].split('.')[0]))

    if not images:
        print("Error: No images were generated from the PDF.")
        return

    for img_name in images:
        img_path = os.path.join(temp_dir, img_name)
        slide_layout = prs.slide_layouts[6]  # Blank layout
        slide = prs.slides.add_slide(slide_layout)
        
        # Add image as a full-slide picture
        slide.shapes.add_picture(img_path, 0, 0, width=prs.slide_width, height=prs.slide_height)

    prs.save(pptx_file)
    print(f"\nDone! Presentation saved as: {pptx_file}")

    # 4. Cleanup temporary assets
    shutil.rmtree(temp_dir)
    print(f"Cleaned up temporary directory: {temp_dir}")

if __name__ == "__main__":
    # Check for library dependency
    try:
        import pptx
    except ImportError:
        print("Error: 'python-pptx' is not installed. Run 'pip install python-pptx'")
        sys.exit(1)

    # Run the build process
    build_presentation(dpi=400)
